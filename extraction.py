"""
Orbit — Phase 2: text extraction.

Given an Upload record, produce plain text Groq can analyze. Audio files are
transcribed with Groq's Whisper endpoint (this is the one place Groq touches
raw audio — everything else in this module is local, offline extraction).
"""

import os

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation


class ExtractionError(Exception):
    """Raised when a file's text cannot be extracted."""


def extract_pdf(path: str) -> str:
    reader = PdfReader(path)
    parts = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            parts.append(text)
    text = "\n\n".join(parts).strip()
    if not text:
        raise ExtractionError(
            "No selectable text found in this PDF. It may be a scanned "
            "document — try a text-based PDF or paste the content directly."
        )
    return text


def extract_docx(path: str) -> str:
    doc = DocxDocument(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    text = "\n".join(parts).strip()
    if not text:
        raise ExtractionError("This document appears to be empty.")
    return text


def extract_pptx(path: str) -> str:
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        slide_lines.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        slide_lines.append(" | ".join(cells))
        if len(slide_lines) > 1:
            parts.append("\n".join(slide_lines))
    text = "\n\n".join(parts).strip()
    if not text:
        raise ExtractionError("No text content found in this presentation.")
    return text


def extract_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read().strip()
    if not text:
        raise ExtractionError("This text file is empty.")
    return text


def extract_audio(path: str, groq_client) -> str:
    """Transcribe audio via Groq's Whisper endpoint."""
    with open(path, "rb") as f:
        transcription = groq_client.audio.transcriptions.create(
            file=(os.path.basename(path), f.read()),
            model="whisper-large-v3-turbo",
            response_format="text",
        )
    text = str(transcription).strip()
    if not text:
        raise ExtractionError("Could not transcribe any speech from this audio file.")
    return text


def extract_text(upload, upload_folder: str, groq_client=None) -> str:
    """
    Dispatch extraction based on Upload.file_type.
    `upload` is the Upload model instance (or any object with the same
    attributes); `upload_folder` is the base uploads directory.
    """
    if upload.file_type in ("pasted", "youtube"):
        path = os.path.join(upload_folder, str(upload.user_id), upload.stored_filename)
        return extract_txt(path)

    path = os.path.join(upload_folder, str(upload.user_id), upload.stored_filename)

    if upload.file_type == "pdf":
        return extract_pdf(path)
    if upload.file_type == "docx":
        return extract_docx(path)
    if upload.file_type == "pptx":
        return extract_pptx(path)
    if upload.file_type == "txt":
        return extract_txt(path)
    if upload.file_type in ("mp3", "wav", "m4a"):
        if groq_client is None:
            raise ExtractionError("Audio transcription requires the Groq client.")
        return extract_audio(path, groq_client)

    raise ExtractionError(f"Unsupported file type: {upload.file_type}")
